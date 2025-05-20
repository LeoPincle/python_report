import argparse
import os
from pptx import Presentation
from ITSM_Excel.data_loader import load_project_data  # Asumo que existe
from ITSM_PPT.report_generator import create_table_slide  # Asumo que existe
from ITSM_Util.helpers import format_data_for_table  # Asumo que existe

def generate_consolidated_report(projects, output_file="Reporte_Consolidado.pptx"):
    """Genera un PPTX con todos los proyectos combinados."""
    prs = Presentation("template.pptx")
    
    # Slide 0: Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Reporte Consolidado ITSM"
    
    # Agrupar datos de todos los proyectos
    all_change_requests = []
    all_incidents = []
    all_service_requests = []
    
    for project_name, data in projects.items():
        # Añadir prefijo de proyecto a cada entrada
        for cr in data.get("change_requests", []):
            cr["Proyecto"] = project_name
            all_change_requests.append(cr)
        
        for inc in data.get("incidents", []):
            inc["Proyecto"] = project_name
            all_incidents.append(inc)
        
        for sr in data.get("service_requests", []):
            sr["Proyecto"] = project_name
            all_service_requests.append(sr)
    
    # Slide 1: Change Requests consolidados
    if all_change_requests:
        create_table_slide(
            prs,
            title="Change Requests (Todos los Proyectos)",
            data=format_data_for_table(all_change_requests),
            columns=["Proyecto", "ID", "Summary", "Status", "Created"]  # Columnas personalizables
        )
    
    # Slide 2: Incidents consolidados
    if all_incidents:
        create_table_slide(
            prs,
            title="Incidents (Todos los Proyectos)",
            data=format_data_for_table(all_incidents),
            columns=["Proyecto", "ID", "Priority", "Status", "Created"]
        )
    
    # Slide 3: Service Requests consolidados
    if all_service_requests:
        create_table_slide(
            prs,
            title="Service Requests (Todos los Proyectos)",
            data=format_data_for_table(all_service_requests),
            columns=["Proyecto", "ID", "Status", "Created", "Assignee"]
        )
    
    prs.save(output_file)
    print(f"✅ Reporte consolidado guardado como: {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Genera reporte PPTX consolidado.")
    parser.add_argument("--projects", nargs="+", help="Nombres de proyectos (ej: f5 jira)", default=[])
    parser.add_argument("--json_dir", help="Directorio de JSONs", default="data")
    args = parser.parse_args()
    
    # Cargar proyectos (todos o los seleccionados)
    project_names = args.projects if args.projects else [
        f.replace(".json", "") for f in os.listdir(args.json_dir) if f.endswith(".json")
    ]
    
    projects_data = {}
    for name in project_names:
        data = load_project_data(os.path.join(args.json_dir, f"{name}.json"))
        if data:
            projects_data[name] = data
    
    if not projects_data:
        print("⚠️ No se encontraron datos. Verifica los nombres o el directorio.")
    else:
        generate_consolidated_report(projects_data)