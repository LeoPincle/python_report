from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


class ChangeRequestPPT:

    def __init__(self, ppt_template, input_project_name, input_current_week_date_range_file_name):
        self.ppt_template = ppt_template
        self.project_name = input_project_name
        self.get_current_week = input_current_week_date_range_file_name

    # take slide_id and shape_id as parameter & return the required shape
    def get_shape(self, slide_id, shape_id):
        obj_slide = self.ppt_template.slides.get(slide_id)
        obj_shape = None
        for shape in obj_slide.shapes:
            if shape.shape_id == shape_id:
                obj_shape = shape
        return obj_shape

    def fill_data(self):
        date_range = self.get_current_week
        cr_excel_path = "Output//GCO - " + self.project_name + " Weekly ITSM Report - CR - " + date_range + ".xlsx"
        cr_wb = load_workbook(filename=cr_excel_path)
        cr_summary_ws = cr_wb["Overall CR Summary"]
        cr_by_cat_subcat_ws = cr_wb["CR by Cat and Subcat"]
        top5_cr_ws = cr_wb["Top 5 CR Types"]

        """CR Summary Table"""

        cr_summary_title = self.get_shape(8672, 4)
        cr_summary_title.text_frame.clear()
        cr_summary_title.text_frame.text = "Overall Change Request Summary - " + self.project_name

        cr_summary_table = self.get_shape(8672, 8)
        for row in range(2, cr_summary_ws.max_row + 1):
            for column in range(1, cr_summary_ws.max_column + 1):
                cell_data = cr_summary_ws.cell(row=row, column=column).value
                if (str(cell_data) == "0" or cell_data is None) and row != cr_summary_ws.max_row:
                    continue
                else:
                    cr_summary_table.table.cell(row - 1, column - 1).text = str(cell_data)
                cr_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                cr_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = False

                if row == cr_summary_ws.max_row:
                    cr_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = True


        # Alignment
        for row in range(1, len(cr_summary_table.table.rows)):
            for column in range(1, 9):
                cr_summary_table.table.cell(row, column).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        for row in range(1, len(cr_summary_table.table.rows)):
            cr_summary_table.table.cell(row, 9).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT


        """CR by Category and Subcategory Table"""

        cr_by_cat_subcat_title = self.get_shape(8672, 7)
        cr_by_cat_subcat_title.text_frame.clear()
        cr_by_cat_subcat_title.text_frame.text = "CR by Category and Subcategory - " + self.project_name

        cr_by_cat_subcat_table = self.get_shape(8672, 6)

        try:
            for row in range(2, cr_by_cat_subcat_ws.max_row + 1):
                for column in range(1, cr_by_cat_subcat_ws.max_column + 1):
                    cell_data = cr_by_cat_subcat_ws.cell(row=row, column=column).value
                    if (str(cell_data) == "None" or cell_data is None) and row == cr_by_cat_subcat_ws.max_row:
                        continue
                    else:
                        cr_by_cat_subcat_table.table.cell(row - 1, column - 1).text = str(cell_data)

                    cr_by_cat_subcat_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                    if column == 3:
                        cr_by_cat_subcat_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    if row == cr_by_cat_subcat_ws.max_row:
                        cr_by_cat_subcat_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = True

        except IndexError as ie:
            print(ie, " - ", "Unexpected error in 'CR by Category & Subcategory table' in PPT. Please enter data manually by referring to Excel output.\n")


        """Top 5 CR Types Table"""

        top_5_cr_title = self.get_shape(8681, 4)
        top_5_cr_title.text_frame.clear()
        top_5_cr_title.text_frame.text = "Top 5 CR Types - " + self.project_name

        top5_cr_table = self.get_shape(8681, 3)
        for row in range(2, top5_cr_ws.max_row + 1):
            for column in range(1, 3):
                cell_data = top5_cr_ws.cell(row=row, column=column).value
                top5_cr_table.table.cell(row - 1, column - 1).text = str(cell_data)
                top5_cr_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                if column == 2:
                    top5_cr_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
