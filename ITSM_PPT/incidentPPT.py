from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


class IncidentPPT:
    def __init__(self, ppt_template, input_project_name, input_get_current_week_date_range_file_name, input_get_current_week_date_range_full):
        self.ppt_template = ppt_template
        self.get_current_week = input_get_current_week_date_range_file_name
        self.project_name = input_project_name
        self.get_current_week_full = input_get_current_week_date_range_full

    # take slide_id and shape_id as parameter & return the required shape
    def get_shape(self, slide_id, shape_id):
        obj_slide = self.ppt_template.slides.get(slide_id)
        obj_shape = None
        for shape in obj_slide.shapes:
            if shape.shape_id == shape_id:
                obj_shape = shape
        return obj_shape

    def fill_data(self):
        # Load data from Excel output
        date_range = self.get_current_week
        incident_excel_path = "Output//GCO - " + self.project_name + " Weekly ITSM Report - Incident - " + date_range + ".xlsx"
        incident_wb = load_workbook(filename=incident_excel_path)
        incident_summary_ws = incident_wb["Overall Incident Summary"]
        incident_pie_chart_ws = incident_wb["Pie Chart"]
        inc_by_cat_subcat_ws = incident_wb["Incidents by Cat and Subcat"]
        top_5_inc_types_ws = incident_wb["Top 5 Incident Types"]


        """Intro slide data"""
        project_name_textbox = self.get_shape(8679, 6)
        project_name_textbox.text_frame.clear()
        project_name_textbox.text_frame.text = "GCO - " + self.project_name + " - Weekly ITSM Report"

        date_range_string = self.get_current_week_full
        date_range_textbox = self.get_shape(8679, 5)
        date_range_textbox.text_frame.clear()
        date_range_textbox.text_frame.text = date_range_string
        date_range_textbox.text_frame.paragraphs[0].font.size = Pt(12)


        """Incident Summary table"""

        incident_summary_title = self.get_shape(552, 3)
        incident_summary_title.text_frame.clear()
        incident_summary_title.text_frame.text = "Overall Incident Summary - " + self.project_name

        incident_summary_table = self.get_shape(552, 5)
        for row in range(2, incident_summary_ws.max_row + 1):
            for column in range(2, incident_summary_ws.max_column + 1):
                cell_data = incident_summary_ws.cell(row=row, column=column).value
                if (str(cell_data) == "0" or cell_data is None) and row != incident_summary_ws.max_row:
                    continue
                else:
                    incident_summary_table.table.cell(row - 1, column - 1).text = str(cell_data)
                incident_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                incident_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[
                    0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                incident_summary_table.table.cell(row - 1, column - 1).vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                if row == incident_summary_ws.max_row:
                    incident_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = True

        # formatting data to fill Reds for Response and Resolution less than 100%
        for row in range(1, len(incident_summary_table.table.rows) - 2):
            cell_data1 = incident_summary_table.table.cell(row, 5).text
            if cell_data1 != '100%' and cell_data1 is not None and cell_data1 != '':
                incident_summary_table.table.cell(row, 5).fill.solid()
                incident_summary_table.table.cell(row, 5).fill.fore_color.rgb = RGBColor(255, 0, 0)
                incident_summary_table.table.cell(row, 5).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255,
                                                                                                             255)

            cell_data2 = incident_summary_table.table.cell(row, 7).text
            if cell_data2 != '100%' and cell_data2 is not None and cell_data2 != '':
                incident_summary_table.table.cell(row, 7).fill.solid()
                incident_summary_table.table.cell(row, 7).fill.fore_color.rgb = RGBColor(255, 0, 0)
                incident_summary_table.table.cell(row, 7).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255,
                                                                                                             255)


        """Incident Pie chart"""
        incident_pie_chart = self.get_shape(552, 6).chart
        # Categories - Priorities
        categories = [incident_pie_chart_ws.cell(1, 1).value,
                      incident_pie_chart_ws.cell(1, 2).value,
                      incident_pie_chart_ws.cell(1, 3).value,
                      incident_pie_chart_ws.cell(1, 4).value,
                      incident_pie_chart_ws.cell(1, 5).value]
        # Data - Incidents per priority
        data = (incident_pie_chart_ws.cell(2, 1).value,
                incident_pie_chart_ws.cell(2, 2).value,
                incident_pie_chart_ws.cell(2, 3).value,
                incident_pie_chart_ws.cell(2, 4).value,
                incident_pie_chart_ws.cell(2, 5).value)
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('', data)

        incident_pie_chart.replace_data(chart_data)


        """Incident by Category and Subcategory Table"""

        inc_by_cat_subcat_title = self.get_shape(8673, 3)
        inc_by_cat_subcat_title.text_frame.clear()
        inc_by_cat_subcat_title.text_frame.text = "Incidents by Category and Subcategory - " + self.project_name

        inc_by_cat_subcat_table = self.get_shape(8673, 7)

        try:

            for row in range(2, inc_by_cat_subcat_ws.max_row + 1):
                for column in range(1, inc_by_cat_subcat_ws.max_column + 1):
                    cell_data = inc_by_cat_subcat_ws.cell(row=row, column=column).value
                    cell_data = inc_by_cat_subcat_ws.cell(row=row, column=column).value
                    if (str(cell_data) == "0" or cell_data is None) and row != inc_by_cat_subcat_ws.max_row:
                        continue
                    elif (str(cell_data) == "None" or cell_data is None) and row == inc_by_cat_subcat_ws.max_row:
                        continue
                    else:
                        inc_by_cat_subcat_table.table.cell(row - 1, column - 1).text = str(cell_data)
                    inc_by_cat_subcat_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                    inc_by_cat_subcat_table.table.cell(row - 1, column - 1).vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

                    if row == inc_by_cat_subcat_ws.max_row:
                        inc_by_cat_subcat_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = True

            # Alignment
            for row in range(1, len(inc_by_cat_subcat_table.table.rows)):
                for column in range(2, 7):
                    inc_by_cat_subcat_table.table.cell(row, column).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

            for row in range(1, len(inc_by_cat_subcat_table.table.rows)):
                inc_by_cat_subcat_table.table.cell(row, 7).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        except IndexError as ie:
            print(ie, " - ", "Unexpected error in 'Incident by Category & Subcategory table' in PPT. Please enter data manually by referring to Excel output.\n")


        """Top 5 Incident Types"""

        top_5_inc_types_title = self.get_shape(8673, 9)
        top_5_inc_types_title.text_frame.clear()
        top_5_inc_types_title.text_frame.text = "Top 5 Incident Types - " + self.project_name


        top_5_inc_types_table = self.get_shape(8673, 6)
        for row in range(2, top_5_inc_types_ws.max_row + 1):
            for column in range(1, 3):
                cell_data = top_5_inc_types_ws.cell(row=row, column=column).value
                top_5_inc_types_table.table.cell(row - 1, column - 1).text = str(cell_data)
                top_5_inc_types_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                if column == 2:
                    top_5_inc_types_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

