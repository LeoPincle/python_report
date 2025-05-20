from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


class RequestPPT:
    def __init__(self, ppt_template, input_project_name, input_get_current_week_date_range_file_name):
        self.ppt_template = ppt_template
        self.project_name = input_project_name
        self.get_current_week = input_get_current_week_date_range_file_name

    # take slide_id and shape_id as parameter & return the required shape
    def get_shape(self, slide_id, shape_id):
        obj_slide = self.ppt_template.slides.get(slide_id)
        obj_shape = None
        for shape in obj_slide.shapes:
            if shape.shape_id == shape_id:
                obj_shape = shape
        return obj_shape

    def fill_data(self):
        date_range = self.get_current_week
        request_excel_path = "Output//GCO - " + self.project_name + " Weekly ITSM Report - Request - " + date_range + ".xlsx"
        request_wb = load_workbook(filename=request_excel_path)
        request_summary_ws = request_wb["Overall Request Summary"]

        request_item_title = self.get_shape(2186, 5)
        request_item_title.text_frame.clear()
        request_item_title.text_frame.text = "Overall Request Summary - " + self.project_name

        request_summary_table = self.get_shape(2186, 4)

        try:
            for row in range(2, request_summary_ws.max_row + 1):
                for column in range(1, request_summary_ws.max_column + 1):
                    cell_data = request_summary_ws.cell(row=row, column=column).value
                    if (str(cell_data) == "None" or cell_data is None) and row == request_summary_ws.max_row:
                        continue
                    else:
                        request_summary_table.table.cell(row - 1, column - 1).text = str(cell_data)

                    request_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.size = Pt(11)
                    request_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = False
                    if column == 3:
                        request_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    if row == request_summary_ws.max_row:
                        request_summary_table.table.cell(row - 1, column - 1).text_frame.paragraphs[0].font.bold = True

        except IndexError as ie:
            print(ie, " - ", "Unexpected error in 'Request Summary table' in PPT. Please enter data manually by referring to Excel output.\n")
