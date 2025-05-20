from pptx import Presentation

ppt = Presentation("..\\Template\\GCO - HP Weekly ITSM Report - Template.pptx")

i = 1
for slide in ppt.slides:
    j = 1
    print(i, "\tId:", slide.slide_id)
    for shape in slide.shapes:
        print("\tId:", shape.shape_id, "\tName:", shape.name)
        j += 1
    i += 1